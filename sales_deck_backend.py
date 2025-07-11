import os
import json
import uuid
import base64
import logging
from datetime import datetime, timedelta
from typing import Dict, List, Optional
from io import BytesIO
import tempfile
import threading
import time

from flask import Flask, request, jsonify, send_file, abort
from flask_cors import CORS
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from openai import OpenAI
from werkzeug.utils import secure_filename
import requests

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = Flask(__name__)
CORS(app)

# Configuration
OPENAI_API_KEY = os.getenv('OPENAI_API_KEY')
if not OPENAI_API_KEY:
    raise ValueError("OPENAI_API_KEY environment variable is required")

# Initialize OpenAI client (new v1+ syntax)
# Be explicit about parameters to avoid proxy issues
client = None

def get_openai_client():
    """Initialize OpenAI client with explicit parameters"""
    global client
    if client is None:
        try:
            client = OpenAI(
                api_key=OPENAI_API_KEY,
                timeout=30.0,
                max_retries=2
            )
        except Exception as e:
            logger.error(f"Failed to initialize OpenAI client: {e}")
            raise
    return client

# File storage configuration
UPLOAD_FOLDER = 'temp_files'
ALLOWED_EXTENSIONS = {'png', 'jpg', 'jpeg', 'gif'}
MAX_CONTENT_LENGTH = 16 * 1024 * 1024  # 16MB max file size

os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = MAX_CONTENT_LENGTH

# In-memory storage for generated files (cleanup after 1 hour)
generated_files = {}

def cleanup_old_files():
    """Background task to clean up old files"""
    while True:
        try:
            current_time = datetime.now()
            files_to_remove = []
            
            for file_id, file_info in generated_files.items():
                if current_time - file_info['created_at'] > timedelta(hours=1):
                    files_to_remove.append(file_id)
                    if os.path.exists(file_info['path']):
                        os.remove(file_info['path'])
            
            for file_id in files_to_remove:
                del generated_files[file_id]
                
        except Exception as e:
            logger.error(f"Error in cleanup task: {e}")
            
        time.sleep(300)  # Run every 5 minutes

# Start cleanup thread
cleanup_thread = threading.Thread(target=cleanup_old_files, daemon=True)
cleanup_thread.start()

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def validate_input(data):
    """Validate required input fields"""
    required_fields = ['company_name', 'industry', 'buyer_persona', 'main_pain_point', 'use_case']
    
    for field in required_fields:
        if field not in data or not data[field].strip():
            return False, f"Missing or empty required field: {field}"
    
    # Validate field lengths
    if len(data['company_name']) > 100:
        return False, "Company name too long (max 100 characters)"
    if len(data['industry']) > 100:
        return False, "Industry too long (max 100 characters)"
    if len(data['buyer_persona']) > 200:
        return False, "Buyer persona too long (max 200 characters)"
    if len(data['main_pain_point']) > 500:
        return False, "Main pain point too long (max 500 characters)"
    if len(data['use_case']) > 500:
        return False, "Use case too long (max 500 characters)"
    
    return True, ""

def generate_slide_content(company_name, industry, buyer_persona, main_pain_point, use_case):
    """Generate slide content using GPT-4"""
    prompt = f"""
    Create a professional sales deck for {company_name} targeting {buyer_persona} in the {industry} industry.
    
    Key context:
    - Main pain point: {main_pain_point}
    - Use case: {use_case}
    
    Generate content for exactly 8 slides in JSON format. Each slide should have a "title" and "content" field.
    The content should be concise bullet points (3-5 points per slide).
    
    Slides needed:
    1. Title slide - Company introduction
    2. Problem - Focus on the main pain point
    3. Solution - How we solve the problem
    4. Product Overview - Key features and benefits
    5. Use Case - Specific application for this client
    6. Case Study - Success story (can be hypothetical but realistic)
    7. Pricing - Value proposition and next steps
    8. Call to Action - Clear next steps
    
    Return only valid JSON in this format:
    {{
        "slides": [
            {{"title": "Slide Title", "content": ["Bullet point 1", "Bullet point 2", "Bullet point 3"]}},
            ...
        ]
    }}
    """
    
    try:
        # Use new OpenAI client syntax
        client = get_openai_client()
        response = client.chat.completions.create(
            model="gpt-4o-mini",  # Updated model name
            messages=[
                {"role": "system", "content": "You are a professional sales presentation expert. Generate compelling, professional sales deck content."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=2000,
            temperature=0.7
        )
        
        content = response.choices[0].message.content.strip()
        
        # Try to parse JSON
        try:
            return json.loads(content)
        except json.JSONDecodeError:
            # If JSON parsing fails, try to extract JSON from the response
            import re
            json_match = re.search(r'\{.*\}', content, re.DOTALL)
            if json_match:
                return json.loads(json_match.group())
            else:
                raise ValueError("Could not parse JSON from GPT response")
                
    except Exception as e:
        logger.error(f"Error generating content with GPT-4: {e}")
        # Return fallback content if API fails
        return {
            "slides": [
                {
                    "title": f"Welcome to {company_name}",
                    "content": [f"Professional solutions for {industry}", "Driving innovation and growth", "Your trusted partner"]
                },
                {
                    "title": "The Challenge",
                    "content": [main_pain_point, "Industry-wide impact", "Need for solution"]
                },
                {
                    "title": "Our Solution",
                    "content": ["Innovative approach", "Proven methodology", "Measurable results"]
                },
                {
                    "title": "Product Overview",
                    "content": ["Key features", "Benefits", "Competitive advantages"]
                },
                {
                    "title": "Use Case",
                    "content": [use_case, "Implementation strategy", "Expected outcomes"]
                },
                {
                    "title": "Success Story",
                    "content": ["Client challenge", "Our solution", "Results achieved"]
                },
                {
                    "title": "Investment & ROI",
                    "content": ["Competitive pricing", "Clear value proposition", "Return on investment"]
                },
                {
                    "title": "Next Steps",
                    "content": ["Schedule demo", "Discuss implementation", "Begin partnership"]
                }
            ]
        }

def create_powerpoint(slide_data, company_name, logo_path=None):
    """Create PowerPoint presentation"""
    try:
        prs = Presentation()
        
        # Set slide size to widescreen
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        slides_content = slide_data.get('slides', [])
        
        for i, slide_content in enumerate(slides_content):
            # Use title slide layout for first slide, content layout for others
            if i == 0:
                slide_layout = prs.slide_layouts[0]  # Title slide
            else:
                slide_layout = prs.slide_layouts[1]  # Title and content
            
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            title = slide.shapes.title
            title.text = slide_content.get('title', '')
            
            # Style title
            title.text_frame.paragraphs[0].font.size = Pt(32)
            title.text_frame.paragraphs[0].font.bold = True
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
            
            # Add content for non-title slides
            if i > 0 and len(slide.shapes) > 1:
                content_placeholder = slide.shapes.placeholders[1]
                text_frame = content_placeholder.text_frame
                text_frame.clear()
                
                content_items = slide_content.get('content', [])
                for j, item in enumerate(content_items):
                    if j == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    p.text = item
                    p.level = 0
                    p.font.size = Pt(18)
                    p.font.color.rgb = RGBColor(64, 64, 64)
            
            # Add logo to title slide if provided
            if i == 0 and logo_path and os.path.exists(logo_path):
                try:
                    slide.shapes.add_picture(logo_path, Inches(11), Inches(0.5), height=Inches(1))
                except Exception as e:
                    logger.warning(f"Could not add logo: {e}")
        
        return prs
        
    except Exception as e:
        logger.error(f"Error creating PowerPoint: {e}")
        raise

def save_base64_image(base64_data, filename):
    """Save base64 encoded image to file"""
    try:
        # Remove data URL prefix if present
        if ',' in base64_data:
            base64_data = base64_data.split(',')[1]
        
        image_data = base64.b64decode(base64_data)
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        
        with open(filepath, 'wb') as f:
            f.write(image_data)
        
        return filepath
    except Exception as e:
        logger.error(f"Error saving base64 image: {e}")
        return None

@app.route('/health', methods=['GET'])
def health_check():
    """Health check endpoint"""
    return jsonify({"status": "healthy", "timestamp": datetime.now().isoformat()})

@app.route('/generate-deck', methods=['POST'])
def generate_deck():
    """Main endpoint to generate sales deck"""
    try:
        # Parse JSON data
        if request.is_json:
            data = request.get_json()
        else:
            return jsonify({"error": "Content-Type must be application/json"}), 400
        
        logger.info(f"Received request to generate deck for company: {data.get('company_name', 'Unknown')}")
        
        # Validate input
        is_valid, error_msg = validate_input(data)
        if not is_valid:
            return jsonify({"error": error_msg}), 400
        
        # Extract data
        company_name = data['company_name'].strip()
        industry = data['industry'].strip()
        buyer_persona = data['buyer_persona'].strip()
        main_pain_point = data['main_pain_point'].strip()
        use_case = data['use_case'].strip()
        logo_base64 = data.get('logo_base64', '')
        
        # Handle logo upload
        logo_path = None
        if logo_base64:
            logo_filename = f"logo_{uuid.uuid4().hex}.png"
            logo_path = save_base64_image(logo_base64, logo_filename)
        
        # Generate content with GPT-4
        logger.info("Generating slide content with GPT-4...")
        slide_data = generate_slide_content(company_name, industry, buyer_persona, main_pain_point, use_case)
        
        # Create PowerPoint
        logger.info("Creating PowerPoint presentation...")
        prs = create_powerpoint(slide_data, company_name, logo_path)
        
        # Save presentation
        file_id = str(uuid.uuid4())
        filename = f"sales_deck_{company_name.replace(' ', '_')}_{file_id}.pptx"
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        
        prs.save(filepath)
        
        # Store file info
        generated_files[file_id] = {
            'path': filepath,
            'filename': filename,
            'created_at': datetime.now()
        }
        
        # Clean up logo file if it was created
        if logo_path and os.path.exists(logo_path):
            os.remove(logo_path)
        
        logger.info(f"Successfully generated deck with file_id: {file_id}")
        
        return jsonify({
            "success": True,
            "file_id": file_id,
            "download_url": f"/download/{file_id}",
            "filename": filename,
            "slides_generated": len(slide_data.get('slides', [])),
            "expires_at": (datetime.now() + timedelta(hours=1)).isoformat()
        })
        
    except Exception as e:
        logger.error(f"Error generating deck: {e}")
        return jsonify({"error": f"Internal server error: {str(e)}"}), 500

@app.route('/download/<file_id>', methods=['GET'])
def download_file(file_id):
    """Download generated PowerPoint file"""
    try:
        if file_id not in generated_files:
            abort(404)
        
        file_info = generated_files[file_id]
        
        # Check if file still exists
        if not os.path.exists(file_info['path']):
            del generated_files[file_id]
            abort(404)
        
        # Check if file has expired
        if datetime.now() - file_info['created_at'] > timedelta(hours=1):
            os.remove(file_info['path'])
            del generated_files[file_id]
            abort(410)  # Gone
        
        return send_file(
            file_info['path'],
            as_attachment=True,
            download_name=file_info['filename'],
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
        
    except Exception as e:
        logger.error(f"Error downloading file {file_id}: {e}")
        abort(500)

@app.errorhandler(404)
def not_found(error):
    return jsonify({"error": "File not found or expired"}), 404

@app.errorhandler(500)
def internal_error(error):
    return jsonify({"error": "Internal server error"}), 500

if __name__ == '__main__':
    print("Starting Sales Deck Generator Backend...")
    print("Make sure to set OPENAI_API_KEY environment variable")
    print("API will be available at http://localhost:5000")
    print("Health check: GET /health")
    print("Generate deck: POST /generate-deck")
    print("Download: GET /download/<file_id>")
    
    app.run(debug=True, host='0.0.0.0', port=5000)
